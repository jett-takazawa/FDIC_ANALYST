import datetime
import json
import os
import matplotlib.pyplot as plt
import PySimpleGUI as sg
import openai
import re
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
import ast
from openai import OpenAI
import json
import pandas as pd
import json, csv
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx.util import Inches, Pt


BRANCH_DATA = "cagr_dep_data/summed_deposits_with_zip.json"
CORRELATION_TABLE = "cagr_dep_data/correlation_table.csv"

client = OpenAI(
    api_key="_GROK_API_KEY"
)
# === Configuration ===


# Enable interactive mode for independent chart windows

# Enable interactive mode for independent chart windows
plt.ion()

# Load JSON data
with open('cagr_dep_data/summed_deposits_with_zip.json', 'r') as f:
    data = json.load(f)

# === Static data & prompts ===
place_to_zip = {
    "hana":        "96713",
    "kahului":     "96732",
    "kihei":       "96753",
    "lahaina":     "96761",
    "makawao":     "96768",
    "paia":        "96779",
    "wailuku":     "96793",
    "aiea":        "96701",
    "ewa beach":   "96706",
    "kapolei":     "96707",
    "haleiwa":     "96712",
    "kailua":      "96734",
    "kaneohe":     "96744",
    "pearl city":  "96782",
    "wahiawa":     "96786",
    "waianae":     "96792",
    "waipahu":     "96797",
    "downtown":    "96813",
    "ala moana":   "96814",
    "waikiki":     "96814",
    "waialae":     "96816",
    "kahala":      "96816",
    "kaimuki":     "96816",
    "chinatown":   "96817",
    "mapunapuna":  "96819",
    "manoa":       "96822",
    "hawaii kai":  "96825",
    "moiliili":    "96826",
    "hilo":        "96720",
    "kailua-kona": "96740",
    "waimea":      "96743",
    "kealakekua":  "96750",
    "aina haina":  "96821",
    "pahala":      "96777",
    "pahoa":       "96778",
    "mililani":    "96789",
    "ala moana":    "96814",
    "oahu": (
        "96701","96706","96707","96709","96712","96717","96730","96731",
        "96734","96744","96759","96762","96782","96786","96789","96791",
        "96792","96795","96797","96801","96802","96803","96804","96805",
        "96806","96807","96808","96809","96810","96811","96812","96813",
        "96814","96815","96816","96817","96819","96820","96821","96822",
        "96825","96823","96824","96826","96836","96837","96838","96839",
        "96840","96841","96843","96844","96846","96847","96848","96849",
        "96850","96853","96854","96857","96858","96859","96860","96861",
        "96863","96898"
    ),
    "maui": (
        "96708","96713","96732","96753","96790","96761","96767","96768",
        "96779","96768","96788","96784","96793"
    ),
    "big island": (
        "96704","96710","96718","96719","96720","96721","96725","96726",
        "96727","96728","96740","96745","96753","96755","96749","96750",
        "96739","96760","96764","96771","96772","96773","96737","96774",
        "96776","96777","96778","96780","96781","96783","96785","96738"
    ),
    "kauai": (
        "96703","96705","96714","96715","96716","96741","96746","96747",
        "96751","96752","96754","96756"
    ),
    "molokai": (
        "96729","96742","96748","96757","96770","96765","96766","96769",
        "96722","96796"
    ),
    "westoahu": ("wainae", 'kapolei','ewa', 'kunia'),
    "Central Oahu": ('mililani','waipahu','wahiawa','pearl city')
}

place_to_zip = {k.lower(): v for k, v in place_to_zip.items()}



data_agent_prompt = (
    """
    SYSTEM:
You are a Bank of Hawaii data analyst. ONLY REPLY IN A CLEANED JSON FORMAT. 
You have access to the following pre‑cleaned datasets for any requested ZIP code:
  • 3‑Year CAGR for each bank  
  • 3‑Year CAGR for the full market  
  • Total deposits by bank  
  • Tax‑variable: %‑difference from median in numeric format ('singleReturns': -3.7 = -3.7%)  
  • Correlation table between each tax variable and total deposits  
  • A bank of zip codes to "area names" (always use area names when reffering to zip code)
  • A dictionary containing the current tax values for all zipcodes from given year  

Your job is to:
  1. Use the provided document of correlated variables
  2. Among those variables, select the two variables whose local %‑difference from the annual median is most pronounced (differed from the mean) MUST BE INCLUEDED IN CORRELATION TABLE. 
  ****AVOID TAX VARIABLES WHERE ZIPCODE VALUE = 0****
  3. Calculate median from that variable and year. Formula - (96822 Median) = (96822 Current val) * (1 + %‑difference from the annual median) **RETURN MEDIAN**
  4. Gather the following metrics for the ZIP code under analysis:
     • Bank‑of‑Hawaii 3‑year CAGR  
     • Market 3‑year CAGR  
     • Bank‑of‑Hawaii total deposits  
     • Total market deposits 

  5. Create Analysis Title. Find corresponding zip code with area name using the bank of zip codes to area names. e.g. ‘(Hawaii Kai *CORRESPONDs WITH ZIP*) Deposit Analysis in {year}’ or e.g. '(Waipahu and Mililani) Deposit Analysis in {year}', or if large numbe of zips use Island name. Oahu, Maui, Big Island." 
 
  6.  Please respond **only** with a single JSON object—no code fences, no Python dict wrappers, no extra text. The object **must** use this exact schema: 
JSON:

{
  "Analysis_Title":              string,  // e.g. "Wailuku Deposit Analysis in 2024"
  "Notable_TaxVar_1":            string,  // e.g. "Investment interest paid amount"
  "Median_1":                    string,  // e.g. "914"
  "Diff_From_Mean_1":            string,  // e.g. "+234.8%"
  "Correlation_Strength_1":      string,  // e.g. "Strong Positive"
  "Notable_TaxVar_2":            string,  // e.g. "Taxable interest amount"
  "Median_2":                    string,  // e.g. "11,144"
  "Diff_From_Mean_2":            string,  // e.g. "+32.4%"
  "Correlation_Strength_2":      string,  // e.g. "Moderate Positive"
  "3yr_CAGR_BankOH":             number,  // e.g. -2.3
  "3yr_CAGR_Market":             string,  // e.g. "-0.3%"
  "Deposit_BankOH":              number,
  "Deposit_Market":              number,
  "Deposit_Insight":             string,  // one-line summary
  "Tax_Insight":                 string   // one-line summary
}

"""

)



# Prompts for GPT chart-vars (colors from default map will override any GPT output)
pie_chart_prompt = (
    "You are an Bank of Hawaii infographic bot. You will receive a dict summarizing total deposits "
    "by bank for specific ZIP codes and year. Return a single Python dict with keys:\n"
    "if no data is found for a specific ZIP you may skip it. "
    """{
        "categories": str,
        "values": float,
        "colors": str,
        "title": str,
    } """
    
    "and no other text."
)



bar_chart_prompt = pie_chart_prompt.replace("pie chart", "bar chart")
scatterplot_prompt = (
    "You are an infographic bot. You will receive a dict summarizing total deposits "
    "by bank for specific ZIP codes and year. Return a single Python dict with keys:\n"
    "  'x', 'y', 'labels', 'colors', 'title' 'one line description of market including total deposit and cagr for bank of hawaii'\n"
    "and no other text."
)

# Default color map
default_color_map = {
    "Bank of Hawaii":       "#659CD2",
    "First Hawaiian Bank":  "#E17474",
    "Central Pacific Bank": "#CDFFFA",
    "American Savings Bank":"#77B277",
    "Other":                "#949494"
}

#build api call for grok analysis


def build_preview_slide(pie_png: str,
                        bar_png: str,
                        pptx_path: str,
                        analysis: dict,
                        deposit_data: dict,
                        cagr_banks: dict,
                        ):
    print(cagr_banks)
    print(deposit_data)
   


    """
    Creates a PPTX slide with:
      - Title from analysis["Analysis_Title"]
      - Full-slide background + chart image
      - Deposit & tax insight text
      - 3×6 summary table (CAGR and Deposits)
    """
    # 1) Setup Presentation & Slide
    if os.path.exists(pptx_path):
        prs = Presentation(pptx_path)     # load your existing deck
    else:
        prs = Presentation()              # first time: new deck
        prs.slide_width  = Inches(16)
        prs.slide_height = Inches(9)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 2) Title
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.1), Inches(9), Inches(1))
    title_tf  = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text     = analysis.get("Analysis_Title", "")
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.underline = True
    p.font.color.rgb = RGBColor(0x4F, 0x81, 0xBD)

    slide.shapes.add_picture(
        "PowerPointPICS/bohlogo.jpg",
        left=Inches(12.7), top=Inches(8.45),
        width=Inches(3.2), height=Inches(0.375)
    )

    slide.shapes.add_picture(
        "PowerPointPICS/bottom left flower.jpg",
        left=Inches(0), top=Inches(4.5),
        width=Inches(4.5), height=Inches(4.5)
    )
    
    # 4) Chart image
    slide.shapes.add_picture(
        bar_png,
        left= Inches(10.5),
        top=Inches(0.1),
        width=Inches(5), height=Inches(4.25)
    )

    slide.shapes.add_picture( 
        pie_png,
        left= Inches(1),
        top=Inches(0.7),
        width=Inches(4.3), height=Inches(4.1)
    )


    # 5) Deposit Insight text

    deposit_txt = (
  
        f"Deposit Insight: {analysis.get('Deposit_Insight', '')}"
    )

    left_depsum = Inches(0.5)
    top_depsum  = Inches(4.25)
    width_depsum  = Inches(15)
    height_depsum = Inches(1.2)

    snip_box = slide.shapes.add_textbox(left_depsum, top_depsum, width_depsum, height_depsum)
    snip_tf  = snip_box.text_frame
    snip_tf.text = deposit_txt
    for para in snip_tf.paragraphs:
        para.font.size = Pt(14)

    # 6) Tax Insight text
    tax_txt = (
        f"Tax Insight: {analysis.get('Tax_Insight','')}"
    )

    left_taxsum = Inches(0.5)
    top_taxsum  = Inches(4.6)
    width_taxsum  = Inches(15)
    height_taxsum = Inches(1.2)
    tax_box = slide.shapes.add_textbox(left_taxsum,top_taxsum,width_taxsum,height_taxsum)
    tax_tf  = tax_box.text_frame
    tax_tf.text = tax_txt
    for para in tax_tf.paragraphs:
        para.font.size = Pt(14)

    # 7) 3×6 Summary Table
    cols = [
        "Company Name:",
        "Bank of Hawaii",
        "First Hawaiian Bank",
        "American Savings Bank",
        "Central Pacific Bank",
        "Other",
        "Market"
    ]


    key_dictionary = {
        "CAGRS": "3-YR-CAGR",
        "Depots": "Deposits"
    }
    
    rows = 3
    cols_count = len(cols)
    left = Inches(0.5)
    top  = Inches(5)
    width  = Inches(15)
    height = Inches(1.2)
    tbl = slide.shapes.add_table(rows, cols_count, left, top, width, height).table

    # column widths
    for col in tbl.columns:
        col.width = Inches(15/cols_count)

    # header row
    for i, name in enumerate(cols):
        tbl.cell(0, i).text = name

    # CAGR row (cast to float before formatting)
    for col_idx, name in enumerate(cols):
        cell = tbl.cell(1, col_idx)

        if name == "Company Name:":
            # in the first column put your label
            cell.text = key_dictionary["CAGRS"]
        else:
            # pick the right source for Market vs each bank
            if name == "Market":
                cell.text = analysis.get("3yr_CAGR_Market", 0)
            else:
                cell.text = cagr_banks.get(name, 0)
    

# Row 2: Deposits
    for i, name in enumerate(cols):
        cell = tbl.cell(2, i)
        if name == "Company Name:":
            cell.text = key_dictionary["Depots"]
        else:
            if name == "Market":
                raw_dep = analysis.get('Deposit_Market', 0)
            
            else:
                raw_dep = deposit_data.get(name, 0)
            try:
                dep = float(raw_dep)
            except (TypeError, ValueError):
                dep = 0.0
            tbl.cell(2, i).text = f"${dep:,.1f}M"
    

    cols_tax = [
        "Tax Info",
        analysis["Notable_TaxVar_1"],
        analysis["Notable_TaxVar_2"]
    ]
    rows_tax = {
        "Tax_var": "Tax Variable",
        "Median": "Median",
        "Diff_From_Med": "Difference From Median",
        "Strength": "Correlation Strength"
    }
    

    rows = 4
    cols_count = len(cols_tax)
    left = Inches(0.5)
    top  = Inches(6.75)
    width  = Inches(15)
    height = Inches(1.5)
    tbl_tax = slide.shapes.add_table(rows, cols_count, left, top, width, height).table

    # column widths
    for col in tbl_tax.columns:
        col.width = Inches(15/cols_count)

    # header row
    for i, name in enumerate(cols_tax):
        tbl_tax.cell(0, i).text = name

    # CAGR row (cast to float before formatting)
    rownumbah = 0
    for col_idx, name in enumerate(cols_tax):
        rownumbah += 1
        cell = tbl_tax.cell(1, col_idx) 

        if name == "Tax Info":
            # in the first column put your label
            cell.text = rows_tax["Median"]
        elif rownumbah == 2:
           cell.text = analysis["Median_1"]
        elif rownumbah == 3:
           cell.text = analysis["Median_2"]
        

# Row 2: Deposits
    rownumbah = 0
    for col_idx, name in enumerate(cols_tax):
        rownumbah += 1
        cell = tbl_tax.cell(2, col_idx) 


        if name == "Tax Info":
            # in the first column put your label
            cell.text = rows_tax["Diff_From_Med"]
        elif rownumbah == 2:
            cell.text = analysis["Diff_From_Mean_1"]
        elif rownumbah == 3:
            cell.text = analysis["Diff_From_Mean_2"]
    
    rownumbah = 0
    for col_idx, name in enumerate(cols_tax):
        rownumbah += 1
        cell = tbl_tax.cell(3, col_idx) 

        if name == "Tax Info":
            # in the first column put your label
            cell.text = rows_tax["Strength"]
        elif rownumbah == 2:
            cell.text = analysis["Correlation_Strength_1"]
        elif rownumbah == 3:
            cell.text = analysis["Correlation_Strength_2"]




    # 8) Save
    prs.save(pptx_path)
    print(f"Preview slide saved to: {pptx_path}")




def load_correlation_table(
    csv_path: str
) -> Tuple[List[Dict[str, object]], Dict[str, Dict[str, object]]]:
    """
    Reads a CSV with columns:
      Variable, Strength, Direction, Correlation

    Returns:
      - sorted_list: List[{
            'Variable': str,
            'Strength': str,
            'Direction': str,
            'Correlation': float
        }] sorted by Correlation desc
      - lookup: {
            Variable: {
              'Strength': str,
              'Direction': str,
              'Correlation': float
            }, ...
        }
    """
    rows: List[Dict[str, object]] = []
    lookup: Dict[str, Dict[str, object]] = {}

    with open(csv_path, newline='') as f:
        reader = csv.DictReader(f)
        for r in reader:
            var = r['Variable']
            corr = float(r['Correlation'])
            entry = {
                'Variable': var,
                'Strength': r['Strength'],
                'Direction': r['Direction'],
                'Correlation': corr
            }
            rows.append(entry)
            lookup[var] = {
                'Strength': r['Strength'],
                'Direction': r['Direction'],
                'Correlation': corr
            }

    # sort descending by correlation
    rows.sort(key=lambda x: x['Correlation'], reverse=True)
    return rows, lookup


def get_data_analysis(
    zipcodes, year,
    bank_cagrs, market_cagr, deviation_table,annual_deposits):
    """
    zipcodes:           ['96822', …]
    year:               2024
    bank_cagrs:         {'96822': 0.052, …}
    market_cagr:        {'96822': 0.037, …}
    deviation_table:    {'96822': {'foreignTaxCreditAmount': 0.14, …}, …}
    correlation_csv_path: 'cagr_dep_data/correlation_table.csv'
    """
    # 1) Load your CSV
    y_int       = int(year)
    data_year   = str(min(y_int, 2022))
    data_medians = pd.DataFrame(json.load(open(f"cagr_dep_data/data_{data_year}.json")))
    dict_by_column = data_medians.to_dict()

    sorted_corrs, corr_lookup = load_correlation_table(
    'cagr_dep_data/correlation_table.csv'
)
    # 2) Build the full payload
    payload = {
        "ZIP_Codes":        zipcodes,
        "Year":             year,
        "Bank_CAGRs":       bank_cagrs,
        "Market_CAGR":      market_cagr,
        "Tax_Deviation":    deviation_table,
        "Correlation_Table": sorted_corrs,
        "Annual_Deposit_Data": annual_deposits,
        "Medians_Data": dict_by_column
    }

    # 3) Call ChatGPT
    resp = client.chat.completions.create(
        model="gpt-4.1-2025-04-14",
        temperature=0,
        messages=[
            {"role":"system", "content": data_agent_prompt},
            {"role":"user", "content":
                "Here are the datasets for analysis:\n\n"
                + json.dumps(payload, indent=2)
            }
        ]
    )

    # 4) Parse JSON (or return raw on failure)
    content = resp.choices[0].message.content
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        return {"raw": content}



def compute_deviation_table(year: str, zip_codes: list[str]) -> dict[str, dict[str, float]]:
    """
    For the given year and list of ZIP codes, load data_{year}.json (capped at 2022),
    compute medians (ignoring zeros), then return:
      {
        "96822": { "investmentInterestPaidAmount": 52.6, ... },
        "96823": { ... },
        ...
      }
    where each number is a plain Python float.
    """
    # 1. Cap at 2022, load DataFrame
    y_int       = int(year)
    data_year   = str(min(y_int, 2022))
    df = pd.DataFrame(json.load(open(f"cagr_dep_data/data_{data_year}.json")))

    # 2. Numeric cols only, coerce to numbers
    numeric = df.drop(columns=['zipCode','year'])
    numeric = numeric.apply(pd.to_numeric, errors='coerce')

    # 3. Compute medians ignoring zeros
    medians = numeric.replace(0, pd.NA).median().astype(float)

    # 4. Build result dict
    deviation_table: dict[str, dict[str, float]] = {}

    for z in zip_codes:
        # select the one-row record for this ZIP/year
        rec = df[
            (df['zipCode'].astype(str) == z) &
            (df['year'] == int(data_year))
        ]
        if rec.empty:
            # no data for this ZIP → skip or fill zeros
            deviation_table[z] = {}
            continue

        # get numeric series for that ZIP
        vals = rec.drop(columns=['zipCode','year']).iloc[0].astype(float)

        # compute % difference
        pct_diff = ((vals - medians) / medians * 100).round(1)

        # keep only Python floats & skip NaN
        deviation_table[z] = {
            category: float(pct_diff.loc[category])
            for category in medians.index
            if pd.notna(pct_diff.loc[category])
        }

    return deviation_table


# Fetch chart variables via OpenAI (GPT)
def get_chart_vars(prompt_text, user_query, summary):
    resp =client.chat.completions.create(
        model="gpt-4.1-2025-04-14",
        temperature=0,
        messages=[
            {"role":"system","content":prompt_text},
            {"role":"user","content":f"{user_query}\nData summary: {summary}"}
        ]
    )
    try:
        return ast.literal_eval(resp.choices[0].message.content.strip())
    except:
        return {}

# Summarize deposits for one or more ZIP codes
def summarize_deposits(zipcodes, year):
    if isinstance(zipcodes, str):
        zipcodes = [zipcodes]
    key = f"Deposits, {year} (in thousands)"
    summary = {bank: 0.0 for bank in default_color_map}
    for rec in data:
        if rec.get('Zip Code') not in zipcodes:
            continue
        amt = rec.get(key, 0) / 1000.0
        matched = False
        for bank in default_color_map:
            if bank != "Other" and bank.upper() in rec.get('Parent Bank', '').upper():
                summary[bank] += amt
                matched = True
                break
        if not matched:
            summary["Other"] += amt
    
    print(summary)
    return summary



# Build GUI
sg.theme("DarkBlue2")
layout = [
    [sg.Text("Enter command:"), sg.InputText(key='-IN-', size=(70,3))],
    [sg.Button('Submit', bind_return_key=True)]
]
window = sg.Window('DataViz Agent', layout, finalize=True)

while True:
    event, vals = window.read()
    if event in (sg.WIN_CLOSED, 'Exit'):
        break

    query = vals['-IN-']
    q_lower = query.lower()

    # Extract ZIPs from numbers or place names
    zip_codes = re.findall(r"\b(\d{5})\b", q_lower)
    if not zip_codes:
        zip_codes = [v for place, v in place_to_zip.items() if place in q_lower]

    if zip_codes and isinstance(zip_codes[0], tuple):
        zip_codes = list(zip_codes[0])

# now zip_codes is a flat list of strings
    print(zip_codes)

    if not zip_codes:
        zip_codes = ['96822']

    # Year
    m = re.search(r"\b(201[2-9]|202[0-4])\b", q_lower)
    year = m.group(1) if m else '2024'
    


    # now proceed with your existing summary/chart code
    # Data summary
    prev_year = str(int(year) - 3)
    summary_prev    = summarize_deposits(zip_codes, prev_year)
    summary_current = summarize_deposits(zip_codes, year)
    deviation_table = compute_deviation_table(year, zip_codes)
    print(deviation_table)


# (optional) compute totals for each
    total_prev    = sum(summary_prev.values())
    total_current = sum(summary_current.values())
    



# Setting up data for data analysis bot
    bank_cagrs = {}
    for bank in summary_current:
        start = summary_prev[bank]
        end   = summary_current[bank]
        if start > 0:
            cagr = ((end / start) ** (1/3) - 1) * 100
            bank_cagrs[bank] = f"{cagr:.1f}%"
        else:
            bank_cagrs[bank] = "N/A"
    bank_cagrs_float = {}
    for bank, val in bank_cagrs.items():
        if isinstance(val, str) and val.endswith('%'):
        # strip the '%' and convert to float
            bank_cagrs_float[bank] = float(val.rstrip('%'))
        else:
        # for non-percent values, set None (or maybe float('nan'))
            bank_cagrs_float[bank] = None

    # market-wide CAGR
    market_cagr = ((total_current / total_prev) ** (1/3) - 1) * 100 \
                if total_prev > 0 else None
   
    print("Bank-of-HI 3yr CAGRs:", bank_cagrs)
    print("Market 3yr CAGR:", f"{market_cagr:.1f}%" if market_cagr is not None else "N/A")


    if total_current == 0:
        sg.popup(f"No data for ZIP(s) {', '.join(zip_codes)} in {year}.")
        continue

    # Inline blurb
    ms_pct = summary_current.get("Bank of Hawaii",0)/total_current*100
    zip_label = ' and '.join(zip_codes)
    blurb = f"We hold {ms_pct:.1f}% of the total deposits in ZIP {zip_label} for {year}."

      # Plot pie chart

    vars_dict    = get_chart_vars(pie_chart_prompt, query, summary_current)
    default_cats = list(summary_current.keys())
    raw_cats     = vars_dict.get('categories', default_cats)
    cats         = [c for c in raw_cats if c in default_cats]
    if not cats:
        cats = default_cats

    # ——— 1) filter out zero‐value slices ———
    filtered = [
        (c, summary_current[c], default_color_map.get(c, '#888888'))
        for c in cats
        if summary_current.get(c, 0) > 0
    ]
    # if everything was zero, fall back to all cats
    if not filtered:
        filtered = [
            (c, summary_current[c], default_color_map.get(c, '#888888'))
            for c in cats
        ]

    filtered_cats, vals, cols = zip(*filtered)
    title = vars_dict.get('title', f"Market Share – {zip_label} ({year})")

    # ——— 2) force a square figure of a fixed size ———
    fig, ax = plt.subplots(figsize=(7, 7), dpi=100)

    # ——— 3) draw the pie with no on-slice labels ———
    ax.pie(
        vals,
        labels=None,
        colors=cols,
        explode=[0.1] + [0] * (len(vals) - 1),
        autopct=lambda pct: f'{pct:.1f}%' if pct > 0 else '',
        startangle=90,
        pctdistance=0.8
    )

    # ——— 4) perfect circle ———
    ax.axis('equal')

    # ——— 5) title styling ———
    ax.set_title(
        title,
        fontdict={'family': 'Arial', 'size': 16, 'weight': 'bold'}
    )

    # ——— 6) clean external legend ———
    ax.legend(
        filtered_cats,
        loc='center left',
        bbox_to_anchor=(1, 0.5),
        frameon=False
    )

    # ——— 7) tighten layout so saved PNG is exactly 7×7 in @ 300 DPI ———
    fig.tight_layout()

    # --- SAVE AS PNG INSTEAD OF SHOWING ---
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_title = title.replace(" ", "_").replace(":", "")
    filename_pie = f"Graphs/piechart_{ts}.png"
    fig.savefig(filename_pie, dpi=300, bbox_inches='tight')

    print(f"📈 Chart saved to {filename_pie}")

# -------------------------------------

   # plot bar chart
    vars_dict = get_chart_vars(bar_chart_prompt, query, summary_current)
    default_cats = list(summary_current.keys())
    raw_cats = vars_dict.get('categories', default_cats)
    cats = [c for c in raw_cats if c in default_cats]
    if not cats:
        cats = default_cats
    vals = [summary_current[c] for c in cats]
    cols = [default_color_map.get(c, '#888888') for c in cats]
    title = vars_dict.get('title', f"Total Deposits – {zip_label} ({year})")
    


    fig, ax = plt.subplots(figsize=(9, 7))
    ax.bar(cats, vals, color=cols)
    ax.set_ylabel('Deposits (Millions USD)')
    ax.set_title(
        title,
        fontdict={
            'family': 'Arial',
            'size': 14,
            'weight': 'bold'
        })
    plt.setp(ax.get_xticklabels(), rotation=45, ha='right')
    fig.subplots_adjust(bottom=0.25)
    
    
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    # make sure your title is filesystem-safe (no slashes or colons)
    safe_title = title.replace(" ", "_").replace(":", "")
    filename_bar = f"Graphs/barchart_{ts}.png"
    #you can also specify a directory here:
    


    fig.savefig(filename_bar, dpi=300, bbox_inches='tight')
    print(f"📈 Chart saved to {filename_bar}")
# -------------------------------------




    #Run data analaytics function 
    
    data_analyzed = get_data_analysis(zip_codes, year,
    bank_cagrs, market_cagr, deviation_table, summary_current)
    print(data_analyzed)
    pptxname = data_analyzed["Analysis_Title"]
    powerpoint_path = f"PowerPoints/{ts}_{pptxname}.pptx"
    

    build_preview_slide(filename_pie,filename_bar, powerpoint_path,data_analyzed,summary_current,bank_cagrs)

window.close()
